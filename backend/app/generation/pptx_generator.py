"""PptxGenerator — generate PPT presentations from Word documents.

Pipeline: validate Word file -> parse -> LLM summarize -> build PPTX.
"""

from __future__ import annotations

import json
import logging
import re
import time
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.generation.word_parser import WordParser
from app.llm.base import BaseLLMProvider

logger = logging.getLogger(__name__)

# Widescreen slide size in inches
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

# Colors
_DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
_MED_BLUE = RGBColor(0x2C, 0x52, 0x82)
_LIGHT_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)
_BLACK = RGBColor(0x00, 0x00, 0x00)


@dataclass
class SlideContent:
    """A single slide's content."""

    slide_type: str  # cover | toc | chapter | conclusion
    title: str
    bullets: list[str] = field(default_factory=list)


@dataclass
class PptxResult:
    """Result of a PPT generation run."""

    output_path: Path
    slide_count: int
    slides: list[SlideContent]
    source_doc: str
    duration_ms: int


class PptxGenerator:
    """Generate PPT presentations from Word documents."""

    def __init__(self, llm: BaseLLMProvider, output_dir: Path) -> None:
        self._llm = llm
        self._output_dir = Path(output_dir)
        self._word_parser = WordParser()

    async def generate(self, file_path: Path, template_path: Path | None = None) -> PptxResult:
        """Full pipeline: validate -> parse Word -> LLM summarize -> build PPTX."""
        start_ms = int(time.time() * 1000)

        file_path = Path(file_path)
        self._word_parser.validate(file_path)
        structure = self._word_parser.parse(file_path)

        slides = await self._summarize_sections(title=structure.title, sections=structure.sections)
        slides = self._ensure_structural_slides(slides, structure.title)

        output = self._build_pptx(slides, title=structure.title, template_path=template_path)

        duration_ms = int(time.time() * 1000) - start_ms

        return PptxResult(
            output_path=output,
            slide_count=len(slides),
            slides=slides,
            source_doc=structure.title,
            duration_ms=duration_ms,
        )

    async def _summarize_sections(self, title: str, sections: list) -> list[SlideContent]:
        """Send all sections to LLM once, get JSON array of slides."""
        section_descriptions = []
        for sec in sections:
            heading = sec.heading if sec.heading else "(无标题)"
            paras_preview = sec.paragraphs[:3]
            section_descriptions.append(f"章节: {heading}\n段落摘要: {'; '.join(paras_preview)}")

        sections_text = "\n\n".join(section_descriptions)
        prompt = (
            f"请将以下文档内容转换为PPT幻灯片。文档标题: {title}\n\n"
            f"文档内容:\n{sections_text}\n\n"
            "请返回JSON格式，包含slides数组。每张幻灯片有slide_type"
            "(cover|toc|chapter|conclusion)、title、bullets字段。\n"
            '格式: {"slides": [{"slide_type": "...", "title": "...", "bullets": [...]}]}\n'
            "确保第一张是cover，第二张是toc，最后一张是conclusion。"
        )

        response = await self._llm.chat([{"role": "user", "content": prompt}])
        return self._parse_llm_response(response)

    def _parse_llm_response(self, response: str) -> list[SlideContent]:
        """Parse LLM response into a list of SlideContent."""
        # Try direct parse first
        try:
            data = json.loads(response)
        except json.JSONDecodeError:
            # Try to extract JSON from the response text
            match = re.search(r"\{.*\}", response, re.DOTALL)
            if not match:
                return []
            try:
                data = json.loads(match.group())
            except json.JSONDecodeError:
                return []

        if not isinstance(data, dict) or "slides" not in data:
            return []

        slides_raw = data["slides"]
        if not isinstance(slides_raw, list):
            return []

        result: list[SlideContent] = []
        for item in slides_raw:
            if not isinstance(item, dict):
                continue
            slide_type = item.get("slide_type", "chapter")
            title = item.get("title", "")
            bullets = item.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = []
            result.append(SlideContent(slide_type=slide_type, title=title, bullets=bullets))
        return result

    def _ensure_structural_slides(self, slides: list[SlideContent], doc_title: str) -> list[SlideContent]:
        """Ensure cover, toc, and conclusion slides are always present."""
        result = list(slides)

        # Ensure cover
        if not result or result[0].slide_type != "cover":
            result.insert(0, SlideContent(slide_type="cover", title=doc_title))

        # Ensure toc after cover
        if len(result) < 2 or result[1].slide_type != "toc":
            # Collect chapter titles for TOC bullets
            toc_bullets = [s.title for s in result if s.slide_type == "chapter" and s.title]
            result.insert(1, SlideContent(slide_type="toc", title="目录", bullets=toc_bullets))

        # Ensure conclusion at end
        if not result or result[-1].slide_type != "conclusion":
            result.append(SlideContent(slide_type="conclusion", title="谢谢", bullets=[]))

        return result

    def _build_pptx(
        self,
        slides: list[SlideContent],
        title: str,
        template_path: Path | None = None,
    ) -> Path:
        """Build PPTX file with python-pptx."""
        if template_path and Path(template_path).exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
            prs.slide_width = _SLIDE_WIDTH
            prs.slide_height = _SLIDE_HEIGHT

        # Use blank layout (index 6 in default template)
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        for slide_data in slides:
            slide_layout = self._find_layout(prs, slide_data.slide_type, blank_layout)
            slide = prs.slides.add_slide(slide_layout)
            self._style_slide(slide, slide_data, prs)

        # Save
        self._output_dir.mkdir(parents=True, exist_ok=True)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[\\/:*?"<>|]', "_", title)
        filename = f"PPT_{safe_title}_{timestamp}.pptx"
        output_path = self._output_dir / filename
        prs.save(str(output_path))

        return output_path

    def _find_layout(self, prs: Presentation, slide_type: str, fallback) -> object:
        """Find an appropriate slide layout for the given slide type."""
        # Try to find layout by name
        for layout in prs.slide_layouts:
            name = layout.name.lower()
            if slide_type == "cover" and "title" in name and "section" not in name:
                return layout
            if slide_type == "toc" and ("section" in name or "toc" in name):
                return layout
        return fallback

    def _style_slide(self, slide, data: SlideContent, prs: Presentation) -> None:
        """Apply styling to a slide based on its type."""
        st = data.slide_type

        if st == "cover":
            self._style_cover(slide, data)
        elif st == "toc":
            self._style_toc(slide, data)
        elif st == "conclusion":
            self._style_conclusion(slide, data)
        else:
            self._style_chapter(slide, data)

    def _style_cover(self, slide, data: SlideContent) -> None:
        """Dark blue gradient background, white title text."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _DARK_BLUE

        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(10.333)
        height = Inches(2.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.title
        p.font.size = Pt(40)
        p.font.color.rgb = _WHITE
        p.font.bold = True
        p.alignment = 1  # PP_ALIGN.CENTER

    def _style_toc(self, slide, data: SlideContent) -> None:
        """Light gray background, blue title, bullet list."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _LIGHT_GRAY

        # Title
        left = Inches(1.0)
        top = Inches(0.5)
        width = Inches(11.333)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = data.title
        p.font.size = Pt(32)
        p.font.color.rgb = _LIGHT_BLUE
        p.font.bold = True

        # Bullets
        if data.bullets:
            left = Inches(1.5)
            top = Inches(2.0)
            width = Inches(10.333)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(data.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = _BLACK
                p.space_after = Pt(8)

    def _style_chapter(self, slide, data: SlideContent) -> None:
        """White background, blue left bar decoration, title + bullets."""
        # Blue left bar
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.15)
        height = _SLIDE_HEIGHT
        shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = _LIGHT_BLUE
        shape.line.fill.background()

        # Title
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.533)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = data.title
        p.font.size = Pt(28)
        p.font.color.rgb = _DARK_BLUE
        p.font.bold = True

        # Bullets
        if data.bullets:
            left = Inches(1.2)
            top = Inches(2.0)
            width = Inches(11.133)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(data.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = _BLACK
                p.space_after = Pt(6)

    def _style_conclusion(self, slide, data: SlideContent) -> None:
        """Same as cover, large text."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _DARK_BLUE

        left = Inches(2.0)
        top = Inches(2.0)
        width = Inches(9.333)
        height = Inches(3.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = data.title or "谢谢"
        p.font.size = Pt(54)
        p.font.color.rgb = _WHITE
        p.font.bold = True
        p.alignment = 1  # PP_ALIGN.CENTER
