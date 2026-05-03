from __future__ import annotations

import logging
from collections.abc import Callable
from pathlib import Path

import subprocess
import tempfile

import pytesseract
from PIL import Image

from app.config import OCRConfig
from app.models.document import ExtractedDoc, FileInfo, StructureItem

logger = logging.getLogger(__name__)


class Extractor:
    def __init__(self, ocr_config: OCRConfig) -> None:
        if ocr_config.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = ocr_config.tesseract_cmd
        self._handlers: dict[str, Callable[[Path], ExtractedDoc]] = {
            ".doc": self._extract_doc,
            ".docx": self._extract_docx,
            ".pdf": self._extract_pdf,
            ".xls": self._extract_xls,
            ".xlsx": self._extract_xlsx,
            ".ppt": self._extract_ppt,
            ".pptx": self._extract_pptx,
            ".png": self._extract_image,
            ".jpg": self._extract_image,
            ".jpeg": self._extract_image,
            ".txt": self._extract_txt,
        }

    def extract(self, file_info: FileInfo) -> ExtractedDoc:
        handler = self._handlers.get(file_info.format)
        if not handler:
            raise ValueError(f"不支持的格式: {file_info.format}")
        return handler(file_info.path)

    def _extract_txt(self, path: Path) -> ExtractedDoc:
        for encoding in ("utf-8", "gb18030"):
            try:
                text = path.read_text(encoding=encoding)
                return ExtractedDoc(text=text, structure=[], source_path=path)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return ExtractedDoc(text="", structure=[], source_path=path)

    def _extract_docx(self, path: Path) -> ExtractedDoc:
        from docx import Document

        doc = Document(str(path))
        text_parts: list[str] = []
        structure: list[StructureItem] = []
        pos = 0

        for para in doc.paragraphs:
            if para.style and para.style.name.startswith("Heading"):
                level_str = para.style.name.replace("Heading ", "").replace("Heading", "1")
                try:
                    level = int(level_str)
                except ValueError:
                    level = 1
                structure.append(StructureItem(level=level, text=para.text, position=pos))
            text_parts.append(para.text)
            pos += len(para.text) + 1

        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                text_parts.append(" | ".join(cells))

        return ExtractedDoc(text="\n".join(text_parts), structure=structure, source_path=path)

    def _extract_pdf(self, path: Path) -> ExtractedDoc:
        import pdfplumber

        text_parts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)

        text = "\n".join(text_parts).strip()
        if len(text) < 50:
            try:
                text = self._ocr_pdf(path)
            except Exception as e:
                logger.warning("PDF OCR 失败 %s: %s", path, e)

        return ExtractedDoc(text=text, structure=[], source_path=path)

    def _ocr_pdf(self, path: Path) -> str:
        from pdf2image import convert_from_path

        images = convert_from_path(str(path))
        texts = []
        for img in images:
            texts.append(pytesseract.image_to_string(img, lang="chi_sim"))
        return "\n".join(texts)

    def _extract_xlsx(self, path: Path) -> ExtractedDoc:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True)
        text_parts: list[str] = []
        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                header = [str(c) if c else "" for c in rows[0]]
                for row in rows[1:]:
                    values = [str(c) if c else "" for c in row]
                    text_parts.append(", ".join(f"{h}: {v}" for h, v in zip(header, values, strict=False) if v))
        wb.close()
        return ExtractedDoc(text="\n".join(text_parts), structure=[], source_path=path)

    def _extract_pptx(self, path: Path) -> ExtractedDoc:
        from pptx import Presentation

        prs = Presentation(str(path))
        text_parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts.append(shape.text_frame.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text_parts.append(slide.notes_slide.notes_text_frame.text)
        return ExtractedDoc(text="\n".join(text_parts), structure=[], source_path=path)

    def _extract_image(self, path: Path) -> ExtractedDoc:
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="chi_sim")
        return ExtractedDoc(text=text, structure=[], source_path=path)

    # ── Legacy binary formats (.doc, .xls, .ppt) ──────────────
    # TODO: 当前 .doc/.ppt 依赖 macOS textutil 做格式转换。
    # 若未来迁移到非 macOS 环境，需替换为纯 Python 库：
    #   - xlrd 已是纯 Python（.xls）
    #   - 候选方案：sharepoint-to-text / pyxtxt / olefile + 自行解析
    # 参考：https://github.com/Horsmann/sharepoint-to-text

    def _extract_doc(self, path: Path) -> ExtractedDoc:
        """Extract text from .doc via macOS textutil → docx → _extract_docx."""
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            subprocess.run(
                ["textutil", "-convert", "docx", "-output", str(tmp_path), str(path)],
                check=True,
                capture_output=True,
            )
            result = self._extract_docx(tmp_path)
            result.source_path = path
            return result
        except Exception:
            logger.exception("textutil .doc 转换失败: %s", path)
            return ExtractedDoc(text="", structure=[], source_path=path)
        finally:
            tmp_path.unlink(missing_ok=True)

    def _extract_xls(self, path: Path) -> ExtractedDoc:
        """Extract text from .xls using xlrd (pure Python)."""
        import xlrd

        wb = xlrd.open_workbook(str(path))
        text_parts: list[str] = []
        for ws in wb.sheets():
            for row_idx in range(ws.nrows):
                cells = [str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols)]
                line = ", ".join(c for c in cells if c and c != "0.0")
                if line:
                    text_parts.append(line)
        return ExtractedDoc(text="\n".join(text_parts), structure=[], source_path=path)

    def _extract_ppt(self, path: Path) -> ExtractedDoc:
        """Extract text from .ppt via macOS textutil → txt."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            subprocess.run(
                ["textutil", "-convert", "txt", "-output", str(tmp_path), str(path)],
                check=True,
                capture_output=True,
            )
            result = self._extract_txt(tmp_path)
            result.source_path = path
            return result
        except Exception:
            logger.exception("textutil .ppt 转换失败: %s", path)
            return ExtractedDoc(text="", structure=[], source_path=path)
        finally:
            tmp_path.unlink(missing_ok=True)
