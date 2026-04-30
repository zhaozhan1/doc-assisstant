from __future__ import annotations

import os
from pathlib import Path

import pytest

os.chdir(Path(__file__).parent.parent)


@pytest.fixture
def fixtures_dir(tmp_path: Path) -> Path:
    """动态生成各格式的测试样本文件。"""
    (tmp_path / "sample.txt").write_text("这是一段测试文本。\n\n第二段落内容。", encoding="utf-8")

    from docx import Document

    doc = Document()
    doc.add_heading("测试标题", level=1)
    doc.add_paragraph("这是正文内容。")
    doc.save(str(tmp_path / "sample.docx"))

    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["姓名", "部门"])
    ws.append(["张三", "办公室"])
    wb.save(str(tmp_path / "sample.xlsx"))

    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "测试PPT标题"
    slide.placeholders[1].text = "测试PPT内容"
    prs.save(str(tmp_path / "sample.pptx"))

    return tmp_path
