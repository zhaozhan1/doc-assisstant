"""Tests for PptxGenerator — PPT generation service."""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from docx import Document

from app.generation.pptx_generator import (
    PptxGenerator,
    PptxResult,
    SlideContent,
)
from app.generation.word_parser import WordParseError
from app.llm.base import BaseLLMProvider

# ---------------------------------------------------------------------------
# Fake LLM provider
# ---------------------------------------------------------------------------


class FakeLLM(BaseLLMProvider):
    """LLM that returns a canned slide JSON response."""

    async def chat(self, messages: list[dict], **kwargs) -> str:
        return json.dumps(
            {
                "slides": [
                    {"slide_type": "cover", "title": "测试文档", "bullets": []},
                    {"slide_type": "chapter", "title": "背景", "bullets": ["要点1", "要点2"]},
                    {"slide_type": "chapter", "title": "方法", "bullets": ["步骤A"]},
                    {"slide_type": "conclusion", "title": "谢谢", "bullets": []},
                ]
            }
        )

    async def chat_stream(self, messages: list[dict], **kwargs):
        yield ""

    async def embed(self, texts: list[str]) -> list[list[float]]:
        return []


class FakeLLMBadJSON(BaseLLMProvider):
    """LLM that returns invalid JSON wrapped in prose."""

    async def chat(self, messages: list[dict], **kwargs) -> str:
        return '这是前言内容。{"slides": [{"slide_type": "cover", "title": "坏JSON", "bullets": []}]}这是后续内容。'

    async def chat_stream(self, messages: list[dict], **kwargs):
        yield ""

    async def embed(self, texts: list[str]) -> list[list[float]]:
        return []


class FakeLLMEmpty(BaseLLMProvider):
    """LLM that returns no usable JSON."""

    async def chat(self, messages: list[dict], **kwargs) -> str:
        return "抱歉，我无法生成内容。"

    async def chat_stream(self, messages: list[dict], **kwargs):
        yield ""

    async def embed(self, texts: list[str]) -> list[list[float]]:
        return []


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _create_docx(path: Path, paragraphs: list[tuple[str | None, str]]) -> Path:
    """Create a test .docx file with paragraphs."""
    doc = Document()
    for style, text in paragraphs:
        doc.add_paragraph(text, style=style)
    doc.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Tests — _parse_llm_response
# ---------------------------------------------------------------------------


class TestParseLlmResponse:
    """Tests for PptxGenerator._parse_llm_response."""

    def test_valid_json(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        raw = '{"slides": [{"slide_type": "cover", "title": "T", "bullets": []}]}'
        result = gen._parse_llm_response(raw)
        assert len(result) == 1
        assert isinstance(result[0], SlideContent)
        assert result[0].slide_type == "cover"
        assert result[0].title == "T"

    def test_json_embedded_in_text(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        raw = '前面有文字 {"slides": [{"slide_type": "chapter", "title": "X", "bullets": ["a"]}]} 后面也有'
        result = gen._parse_llm_response(raw)
        assert len(result) == 1
        assert result[0].slide_type == "chapter"
        assert result[0].bullets == ["a"]

    def test_invalid_json_returns_empty(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        raw = "完全没有JSON的内容"
        result = gen._parse_llm_response(raw)
        assert result == []

    def test_missing_slides_key_returns_empty(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        raw = '{"data": []}'
        result = gen._parse_llm_response(raw)
        assert result == []


# ---------------------------------------------------------------------------
# Tests — _build_pptx
# ---------------------------------------------------------------------------


class TestBuildPptx:
    """Tests for PptxGenerator._build_pptx (default styling)."""

    def test_produces_valid_pptx(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        slides = [
            SlideContent(slide_type="cover", title="测试标题", bullets=[]),
            SlideContent(slide_type="toc", title="目录", bullets=["章节一", "章节二"]),
            SlideContent(slide_type="chapter", title="章节一", bullets=["要点1", "要点2"]),
            SlideContent(slide_type="conclusion", title="谢谢", bullets=[]),
        ]
        output = gen._build_pptx(slides, title="测试标题", template_path=None)

        assert output.exists()
        assert output.suffix == ".pptx"
        assert output.stat().st_size > 0

        # Verify we can open it and count slides
        from pptx import Presentation

        prs = Presentation(str(output))
        assert len(prs.slides) == 4

    def test_output_file_naming(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        slides = [
            SlideContent(slide_type="cover", title="ABC", bullets=[]),
            SlideContent(slide_type="conclusion", title="谢谢", bullets=[]),
        ]
        output = gen._build_pptx(slides, title="ABC", template_path=None)
        assert output.name.startswith("PPT_ABC_")


# ---------------------------------------------------------------------------
# Tests — _summarize_sections
# ---------------------------------------------------------------------------


class TestSummarizeSections:
    """Tests for PptxGenerator._summarize_sections."""

    @pytest.mark.asyncio
    async def test_returns_slide_content_list(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        from app.generation.word_parser import Section

        sections = [
            Section(level=1, heading="背景", paragraphs=["这是背景内容。"]),
            Section(level=1, heading="方法", paragraphs=["这是方法内容。"]),
        ]
        result = await gen._summarize_sections(title="测试文档", sections=sections)

        assert isinstance(result, list)
        assert all(isinstance(s, SlideContent) for s in result)
        assert len(result) > 0


# ---------------------------------------------------------------------------
# Tests — generate (full pipeline)
# ---------------------------------------------------------------------------


class TestGenerate:
    """Tests for PptxGenerator.generate full pipeline."""

    @pytest.mark.asyncio
    async def test_full_pipeline(self, tmp_path: Path) -> None:
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        docx_path = _create_docx(
            tmp_path / "test.docx",
            [
                ("Heading 1", "总体要求"),
                (None, "这是第一段内容。"),
                ("Heading 1", "工作安排"),
                (None, "这是第二段内容。"),
            ],
        )

        result = await gen.generate(docx_path)

        assert isinstance(result, PptxResult)
        assert result.output_path.exists()
        assert result.slide_count > 0
        assert len(result.slides) > 0
        assert result.source_doc != ""
        assert result.duration_ms >= 0

    @pytest.mark.asyncio
    async def test_bad_json_fallback(self, tmp_path: Path) -> None:
        """If LLM returns bad JSON, generator still produces output with fallback slides."""
        gen = PptxGenerator(llm=FakeLLMBadJSON(), output_dir=tmp_path)
        docx_path = _create_docx(
            tmp_path / "bad.docx",
            [
                ("Heading 1", "标题"),
                (None, "内容。"),
            ],
        )

        result = await gen.generate(docx_path)
        assert isinstance(result, PptxResult)
        assert result.output_path.exists()

    @pytest.mark.asyncio
    async def test_empty_llm_response_fallback(self, tmp_path: Path) -> None:
        """If LLM returns nothing usable, generator creates minimal cover + conclusion."""
        gen = PptxGenerator(llm=FakeLLMEmpty(), output_dir=tmp_path)
        docx_path = _create_docx(
            tmp_path / "empty_resp.docx",
            [
                ("Heading 1", "测试"),
                (None, "内容。"),
            ],
        )

        result = await gen.generate(docx_path)
        assert isinstance(result, PptxResult)
        assert result.output_path.exists()
        # Should have at least cover + conclusion
        assert result.slide_count >= 2
        slide_types = [s.slide_type for s in result.slides]
        assert "cover" in slide_types
        assert "conclusion" in slide_types

    @pytest.mark.asyncio
    async def test_cover_toc_conclusion_always_present(self, tmp_path: Path) -> None:
        """Even if LLM omits some structural slides, generator ensures they exist."""
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        docx_path = _create_docx(
            tmp_path / "struct.docx",
            [("Heading 1", "标题"), (None, "内容。")],
        )

        result = await gen.generate(docx_path)
        slide_types = [s.slide_type for s in result.slides]
        assert slide_types[0] == "cover"
        assert "conclusion" in slide_types

    @pytest.mark.asyncio
    async def test_nonexistent_file_raises(self, tmp_path: Path) -> None:
        """generate() should raise for a nonexistent file."""
        gen = PptxGenerator(llm=FakeLLM(), output_dir=tmp_path)
        with pytest.raises(WordParseError):
            await gen.generate(tmp_path / "nonexistent.docx")
